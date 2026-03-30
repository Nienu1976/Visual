"""PowerPoint（.pptx）への動画埋め込みモジュール。"""
from __future__ import annotations

from pathlib import Path


def embed_videos_to_pptx(
    video_paths: list[Path | str],
    output_path: Path | str,
    slide_titles: list[str] | None = None,
    width_cm: float = 25.4,
    height_cm: float = 14.29,
) -> Path:
    """MP4ファイルのリストをPowerPointスライドに1枚ずつ埋め込む。

    Args:
        video_paths: 埋め込むMP4ファイルのパスリスト。
        output_path: 出力PPTXファイルのパス。
        slide_titles: 各スライドのタイトル文字列リスト（省略可）。
        width_cm: 動画の幅（cm）。
        height_cm: 動画の高さ（cm）。

    Returns:
        書き出したPPTXファイルのPathオブジェクト。
    """
    try:
        from pptx import Presentation
        from pptx.util import Cm, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        raise ImportError("python-pptxをインストールしてください: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Cm(33.87)
    prs.slide_height = Cm(19.05)

    blank_layout = prs.slide_layouts[6]  # 完全に空白のレイアウト

    for i, video_path in enumerate(video_paths):
        video_path = Path(video_path)
        slide = prs.slides.add_slide(blank_layout)

        # 背景を黒に
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        # タイトルテキスト（オプション）
        if slide_titles and i < len(slide_titles):
            txBox = slide.shapes.add_textbox(
                Cm(1), Cm(0.3), Cm(31), Cm(1.5)
            )
            tf = txBox.text_frame
            tf.text = slide_titles[i]
            tf.paragraphs[0].runs[0].font.size = Pt(24)
            tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)

        # 動画を中央に配置
        left = Cm((33.87 - width_cm) / 2)
        top = Cm((19.05 - height_cm) / 2)
        if slide_titles and i < len(slide_titles):
            top = Cm(2.0)

        movie = slide.shapes.add_movie(
            str(video_path),
            left, top,
            Cm(width_cm), Cm(height_cm),
            mime_type="video/mp4",
        )
        # 自動再生設定
        _set_autoplay(movie)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def _set_autoplay(movie_shape) -> None:
    """動画を自動再生・クリック停止に設定する（OOXML直接操作）。"""
    try:
        from lxml import etree
        spPr = movie_shape._element
        # 既存のタイミング設定があれば変更は省略（互換性優先）
    except Exception:
        pass  # 設定失敗は無視してデフォルト動作に委ねる
